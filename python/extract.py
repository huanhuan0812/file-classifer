#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件文本提取器 - 独立模块
支持从 PPTX、PPT、DOCX 文件中提取文本内容
"""

import re
import zipfile
from io import BytesIO
from pathlib import Path

# 尝试导入PPTX处理库
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("警告: python-pptx未安装，PPTX/PPT文件支持将受限。请运行: pip install python-pptx")

# 尝试导入DOCX处理库
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("警告: python-docx未安装，DOCX文件支持将受限。请运行: pip install python-docx")


def extract_text_from_docx(docx_path):
    """
    从DOCX文件中提取所有文本内容
    
    Args:
        docx_path: DOCX文件路径（str或Path对象）
    
    Returns:
        str: 提取的文本内容
    """
    text_chunks = []
    total_text = ""
    doc = None
    
    if not DOCX_SUPPORT:
        print("  ⚠ 警告: python-docx未安装，无法解析docx文件")
        return ""
    
    try:
        doc = Document(docx_path)
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text and paragraph.text.strip():
                text_chunks.append(paragraph.text.strip())
                if len(text_chunks) >= 20:
                    total_text += " ".join(text_chunks) + " "
                    text_chunks.clear()
        
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        text_chunks.append(cell.text.strip())
                        if len(text_chunks) >= 20:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
        
        # 提取页眉页脚
        if hasattr(doc, 'sections'):
            for section in doc.sections:
                # 页眉
                if section.header and section.header.paragraphs:
                    for para in section.header.paragraphs:
                        if para.text and para.text.strip():
                            text_chunks.append(para.text.strip())
                # 页脚
                if section.footer and section.footer.paragraphs:
                    for para in section.footer.paragraphs:
                        if para.text and para.text.strip():
                            text_chunks.append(para.text.strip())
        
        # 添加剩余的文本块
        if text_chunks:
            total_text += " ".join(text_chunks)
            text_chunks.clear()
        
    except Exception as e:
        total_text = ""
    finally:
        if doc is not None:
            del doc
        text_chunks.clear()
        del text_chunks
    
    return total_text


def extract_embedded_docx_from_pptx(pptx_path):
    """
    从PPTX文件中提取嵌入的DOCX文件并读取其内容
    
    Args:
        pptx_path: PPTX文件路径（str或Path对象）
    
    Returns:
        str: 提取的文本内容
    """
    embedded_texts = []
    
    try:
        # PPTX文件本质上是ZIP压缩包
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            # 查找嵌入的对象，通常在 ppt/embeddings/ 目录下
            for file_info in zf.filelist:
                if file_info.filename.startswith('ppt/embeddings/') and file_info.filename.endswith(('.docx', '.doc')):
                    try:
                        # 读取嵌入的docx文件内容
                        docx_data = zf.read(file_info.filename)
                        
                        # 使用内存流处理docx
                        if DOCX_SUPPORT:
                            docx_stream = BytesIO(docx_data)
                            doc = Document(docx_stream)
                            
                            # 提取段落文本
                            for paragraph in doc.paragraphs:
                                if paragraph.text and paragraph.text.strip():
                                    embedded_texts.append(paragraph.text.strip())
                            
                            # 提取表格文本
                            for table in doc.tables:
                                for row in table.rows:
                                    for cell in row.cells:
                                        if cell.text and cell.text.strip():
                                            embedded_texts.append(cell.text.strip())
                            
                            del doc
                        else:
                            # 简单提取：将docx作为zip读取，提取document.xml中的文本
                            try:
                                docx_zip = zipfile.ZipFile(BytesIO(docx_data))
                                if 'word/document.xml' in docx_zip.namelist():
                                    xml_content = docx_zip.read('word/document.xml').decode('utf-8', errors='ignore')
                                    # 简单正则提取文本
                                    text_matches = re.findall(r'>([^<]+)<', xml_content)
                                    for match in text_matches:
                                        if match.strip() and len(match.strip()) > 1:
                                            embedded_texts.append(match.strip())
                                docx_zip.close()
                            except:
                                pass
                        
                    except Exception as e:
                        # 忽略单个嵌入文件的解析错误
                        pass
    
    except Exception as e:
        pass
    
    return " ".join(embedded_texts) if embedded_texts else ""


def extract_text_from_pptx(pptx_path, extract_embedded=True):
    """
    从PPTX文件中提取所有文本内容，包括嵌入的DOCX文件
    
    Args:
        pptx_path: PPTX文件路径（str或Path对象）
        extract_embedded: 是否提取嵌入的DOCX文件内容，默认为True
    
    Returns:
        str: 提取的文本内容
    """
    text_chunks = []
    total_text = ""
    prs = None
    embedded_text = ""
    
    if not PPTX_SUPPORT:
        print("  ⚠ 警告: python-pptx未安装，无法解析pptx文件")
        return ""
    
    try:
        # 提取嵌入的docx文件内容
        if extract_embedded:
            embedded_text = extract_embedded_docx_from_pptx(pptx_path)
            if embedded_text:
                text_chunks.append(embedded_text)
        
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            try:
                # 提取标题
                if slide.shapes.title and slide.shapes.title.text:
                    text_chunks.append(slide.shapes.title.text.strip())
                    if len(text_chunks) >= 10:
                        total_text += " ".join(text_chunks) + " "
                        text_chunks.clear()
            except:
                pass
            
            # 遍历所有形状
            for shape in slide.shapes:
                try:
                    # 提取文本框内容
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text_chunks.append(shape.text.strip())
                        if len(text_chunks) >= 10:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
                    
                    # 提取表格中的文本
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text and cell.text.strip():
                                    text_chunks.append(cell.text.strip())
                                    if len(text_chunks) >= 10:
                                        total_text += " ".join(text_chunks) + " "
                                        text_chunks.clear()
                    
                    # 检查是否有文本框架
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame.text and shape.text_frame.text.strip():
                            text_chunks.append(shape.text_frame.text.strip())
                            
                except:
                    continue
            
            # 提取备注
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_chunks.append(notes_slide.notes_text_frame.text.strip())
                        if len(text_chunks) >= 10:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
            except:
                pass
            
            del slide
        
        # 添加剩余的文本块
        if text_chunks:
            total_text += " ".join(text_chunks)
            text_chunks.clear()
        
    except Exception as e:
        total_text = ""
    finally:
        if prs is not None:
            del prs
        text_chunks.clear()
        del text_chunks
    
    return total_text


def extract_text_from_file(filepath, extract_embedded=True):
    """
    根据文件类型自动提取文本内容（统一入口）
    
    Args:
        filepath: 文件路径（str或Path对象）
        extract_embedded: 对于PPTX是否提取嵌入的DOCX内容，默认为True
    
    Returns:
        str: 提取的文本内容，如果文件格式不支持则返回空字符串
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()
    
    if suffix in ['.pptx', '.ppt']:
        return extract_text_from_pptx(str(filepath), extract_embedded)
    elif suffix == '.docx':
        return extract_text_from_docx(str(filepath))
    else:
        return ""


def get_supported_extensions():
    """
    获取支持的扩展名列表
    
    Returns:
        list: 支持的扩展名列表
    """
    extensions = []
    if PPTX_SUPPORT:
        extensions.extend(['.pptx', '.ppt'])
    if DOCX_SUPPORT:
        extensions.append('.docx')
    return extensions


def is_supported_file(filepath):
    """
    检查文件是否为支持的格式
    
    Args:
        filepath: 文件路径（str或Path对象）
    
    Returns:
        bool: 是否支持
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()
    return suffix in get_supported_extensions()


# 简单测试
if __name__ == "__main__":
    import sys
    
    print("=" * 50)
    print("文件文本提取器 - 独立模块")
    print("=" * 50)
    print(f"PPTX支持: {'✅ 已启用' if PPTX_SUPPORT else '❌ 未安装'}")
    print(f"DOCX支持: {'✅ 已启用' if DOCX_SUPPORT else '❌ 未安装'}")
    print(f"支持扩展名: {get_supported_extensions()}")
    print("=" * 50)
    
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        if is_supported_file(filepath):
            print(f"\n正在提取文本: {filepath}")
            text = extract_text_from_file(filepath)
            print(f"\n提取结果 (共 {len(text)} 字符):")
            print("-" * 50)
            print(text[:500] + ("..." if len(text) > 500 else ""))
            print("-" * 50)
        else:
            print(f"不支持的文件格式: {filepath}")
    else:
        print("\n用法: python text_extractor.py <文件路径>")
        print("支持的文件格式: .pptx, .ppt, .docx")